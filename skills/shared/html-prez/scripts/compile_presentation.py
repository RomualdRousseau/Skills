#!/usr/bin/env python3
import sys
import os
import re

def extract_number(filename):
    match = re.search(r'(\d+)', filename)
    return int(match.group(1)) if match else 0

def compile_assets(slides_dir, output_pdf, output_pptx):
    slides_dir_abs = os.path.abspath(slides_dir)
    output_pdf_abs = os.path.abspath(output_pdf)
    output_pptx_abs = os.path.abspath(output_pptx)

    # Ensure output parent directories exist
    os.makedirs(os.path.dirname(output_pdf_abs), exist_ok=True)
    os.makedirs(os.path.dirname(output_pptx_abs), exist_ok=True)

    print(f"Source slides directory: {slides_dir_abs}")
    print(f"Target PDF output: {output_pdf_abs}")
    print(f"Target PPTX output: {output_pptx_abs}")

    # Check for Pillow and python-pptx dependencies
    try:
        from PIL import Image
    except ImportError:
        print("Error: The 'Pillow' package is required for compiling PDFs.")
        print("Please install it using: uv pip install Pillow")
        sys.exit(1)

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("Error: The 'python-pptx' package is required for compiling PPTXs.")
        print("Please install it using: uv pip install python-pptx")
        sys.exit(1)

    # 1. Gather and sort PNG slides
    if not os.path.exists(slides_dir_abs):
        print(f"Error: Directory {slides_dir_abs} not found.")
        sys.exit(1)

    png_files = [f for f in os.listdir(slides_dir_abs) if f.lower().endswith(".png")]
    if not png_files:
        print(f"Error: No PNG slide files found in {slides_dir_abs}")
        sys.exit(1)

    # Sort numerically based on numbers in the filename
    png_files.sort(key=extract_number)
    print(f"Found {len(png_files)} PNG files to compile:")
    for f in png_files:
        print(f"  - {f}")
    print()

    full_png_paths = [os.path.join(slides_dir_abs, f) for f in png_files]

    # 2. Compile into PDF
    print("=== Compiling into PDF ===")
    try:
        images = [Image.open(f) for f in full_png_paths]
        # Convert all to RGB to avoid alpha channel / translucent format conflicts in PDF
        rgb_images = [img.convert('RGB') for img in images]
        
        rgb_images[0].save(
            output_pdf_abs,
            format="PDF",
            save_all=True,
            append_images=rgb_images[1:]
        )
        print(f"✓ Saved PDF: {output_pdf_abs}\n")
    except Exception as e:
        print(f"Error compiling PDF: {e}")
        sys.exit(1)

    # 3. Compile into PPTX (Widescreen 16:9)
    print("=== Compiling into PPTX (16:9 Widescreen) ===")
    try:
        prs = Presentation()
        # Set widescreen slide dimensions (16:9 proportion)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        blank_slide_layout = prs.slide_layouts[6] # Blank layout
        
        for img_path in full_png_paths:
            # Add slide
            slide = prs.slides.add_slide(blank_slide_layout)
            # Add picture covering the entire widescreen slide canvas
            slide.shapes.add_picture(
                img_path,
                Inches(0),
                Inches(0),
                width=Inches(13.333),
                height=Inches(7.5)
            )
            
        prs.save(output_pptx_abs)
        print(f"✓ Saved PPTX: {output_pptx_abs}\n")
    except Exception as e:
        print(f"Error compiling PPTX: {e}")
        sys.exit(1)

    print("=== PPTX AND PDF COMPILATIONS COMPLETED SUCCESSFULLY ===")

def main():
    if len(sys.argv) < 4:
        print("Usage: compile_presentation.py <slides_directory> <output_pdf> <output_pptx>")
        print("Example: uv run python3 compile_presentation.py slides_dir/ presentation.pdf presentation.pptx")
        sys.exit(1)

    slides_dir = sys.argv[1]
    output_pdf = sys.argv[2]
    output_pptx = sys.argv[3]

    compile_assets(slides_dir, output_pdf, output_pptx)

if __name__ == "__main__":
    main()
