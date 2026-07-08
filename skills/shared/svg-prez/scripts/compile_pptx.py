#!/usr/bin/env python3
import os
import glob
import sys
import copy
from pptx.opc.package import _Relationship

# Inline Metadata for UV to automatically fetch dependencies
# /// script
# dependencies = [
#   "python-pptx",
#   "svg2pptx",
# ]
# ///

def inject_logo_to_slide(slide, idx, logo_dir):
    """
    Natively injects the Servier logo or symbol into a slide at its exact 
    pixel coordinates (divided by 144.0 to convert to PowerPoint Inches).
    """
    from pptx.util import Inches
    
    if idx == 19: # Closing Slide (Dark Blue BG -> White Symbol)
        logo_path = os.path.join(logo_dir, "Servier_Symbole_Blanc.png")
        if os.path.exists(logo_path):
            slide.shapes.add_picture(
                logo_path,
                Inches(824 / 144.0),
                Inches(404 / 144.0),
                width=Inches(272 / 144.0),
                height=Inches(272 / 144.0)
            )
    elif idx == 1: # Title Slide (Dark Blue BG -> White Signature Logo)
        logo_path = os.path.join(logo_dir, "Servier_Logo_Sign_Blanc.png")
        if os.path.exists(logo_path):
            slide.shapes.add_picture(
                logo_path,
                Inches(1560 / 144.0),
                Inches(960 / 144.0),
                width=Inches(288 / 144.0),
                height=Inches(60 / 144.0)
            )
    elif idx in [2, 5, 9, 13]: # Chapter Slides (White BG -> Colored Signature Logo)
        logo_path = os.path.join(logo_dir, "Servier_Logo_Sign_RVB.png")
        if os.path.exists(logo_path):
            slide.shapes.add_picture(
                logo_path,
                Inches(1560 / 144.0),
                Inches(960 / 144.0),
                width=Inches(288 / 144.0),
                height=Inches(60 / 144.0)
            )
    else: # Standard Content Slides (White BG -> Colored Logo)
        logo_path = os.path.join(logo_dir, "Servier_Logo_RVB.png")
        if os.path.exists(logo_path):
            slide.shapes.add_picture(
                logo_path,
                Inches(1655 / 144.0),
                Inches(1002 / 144.0),
                width=Inches(193 / 144.0),
                height=Inches(60 / 144.0)
            )

def main():
    if len(sys.argv) < 3:
        print("Usage: compile_pptx.py <slides_dir> <output_pptx>")
        sys.exit(1)

    slides_dir = sys.argv[1]
    output_pptx = sys.argv[2]
    
    # Locate brand logos relative to the input folder or workspace
    logo_dir = os.path.dirname(os.path.abspath(output_pptx))

    print("=== NATIVE VECTOR PPTX COMPILER FOR HAND-CRAFTED SLIDES ===")
    print(f"Slides Folder: {slides_dir}")
    print(f"Destination Deck: {output_pptx}")
    print(f"Logo assets directory: {logo_dir}")

    try:
        from svg2pptx import svg_to_pptx
    except ImportError:
        print("Error: The 'svg2pptx' package is required.")
        sys.exit(1)

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("Error: The 'python-pptx' package is required.")
        sys.exit(1)

    # 1. Convert all SVG slides to PPTX slide-by-slide
    print("\n=== Step 1: Converting Hand-Crafted SVGs to PPTXs ===")
    svgs = glob.glob(os.path.join(slides_dir, "*.svg"))
    if not svgs:
        print("Error: No SVG slides found inside the slides folder.")
        sys.exit(1)

    # Sort numerically (Slide1, Slide2, ... Slide19)
    svgs = sorted(svgs, key=lambda x: int(os.path.basename(x).replace("Slide", "").replace(".svg", "")))

    pptx_slides = []
    for svg_path in svgs:
        base_name = os.path.basename(svg_path).replace(".svg", "")
        pptx_path = os.path.join(slides_dir, f"{base_name}.pptx")
        print(f"Converting {base_name}.svg -> {base_name}.pptx...")
        try:
            # 1.1 Convert SVG to raw vector PPTX
            svg_to_pptx(svg_path, pptx_path)
            
            # 1.2 Open the slide and inject the logo natively
            slide_idx = int(base_name.replace("Slide", ""))
            slide_prs = Presentation(pptx_path)
            inject_logo_to_slide(slide_prs.slides[0], slide_idx, logo_dir)
            slide_prs.save(pptx_path)
            
            pptx_slides.append(pptx_path)
        except Exception as e:
            print(f"✗ Failed converting {base_name}.svg: {e}")
            sys.exit(1)

    # 2. Assemble all individual PPTX slides into a single widescreen presentation
    print("\n=== Step 2: Assembling Individual PPTXs into Unified Presentation ===")
    main_prs = Presentation()
    main_prs.slide_width = Inches(13.333)
    main_prs.slide_height = Inches(7.5)

    try:
        blank_layout = main_prs.slide_layouts[6]
    except IndexError:
        blank_layout = main_prs.slide_layouts[0]

    for i, slide_pptx in enumerate(pptx_slides, start=1):
        print(f"Merging Slide {i} / {len(pptx_slides)}...")
        indiv_prs = Presentation(slide_pptx)
        new_slide = main_prs.slides.add_slide(blank_layout)
        source_slide = indiv_prs.slides[0]
        
        # Copy shapes at the XML layer (retaining all complex custom geometries, paths, groups, and freeforms)
        for shape in source_slide.shapes:
            el = shape.element
            new_el = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
            
        # Copy relationships safely (crucial for embedded base64 images and logos)
        for rel in source_slide.part.rels.values():
            # Skip layout and notes relationships to prevent ZIP file duplicates and slide corruption
            if "slideLayout" in rel.reltype or "notesSlide" in rel.reltype:
                continue
                
            # Copy image relationships by importing their raw blobs into the target package
            if "image" in rel.reltype:
                import io
                image_stream = io.BytesIO(rel.target_part.blob)
                target_image_part = main_prs.part.package.get_or_add_image_part(image_stream)
                
                new_slide.part.rels._rels[rel.rId] = _Relationship(
                    new_slide.part.rels._base_uri,
                    rel.rId,
                    rel.reltype,
                    target_mode=rel._target_mode,
                    target=target_image_part,
                )

    # Delete the empty initial slide
    if len(main_prs.slides) > len(pptx_slides):
        id_list = main_prs.slides._sldIdLst
        del id_list[0]

    # Save output
    main_prs.save(output_pptx)
    print(f"\n✓ Native Widescreen Hand-crafted Vector PPTX successfully saved to {output_pptx}")
    print("=== COMPILATION COMPLETE ===")

if __name__ == "__main__":
    main()
